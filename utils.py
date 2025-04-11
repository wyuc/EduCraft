import os
import base64
import re
import logging
from pptx import Presentation
import fitz
from pathlib import Path
from config import BASE_DIR

# Configure logging
logger = logging.getLogger(__name__)


def get_temp_dir(ppt_path):
    return BASE_DIR / 'buffer' / Path(ppt_path).stem

def get_images_dir(ppt_path):
    return get_temp_dir(ppt_path) / 'images'

def preprocess_ppt(ppt_path):
    text_content = extract_text_from_ppt(ppt_path)
    temp_dir = get_temp_dir(ppt_path)
    os.makedirs(temp_dir, exist_ok=True)

    
    pptx_to_pdf(ppt_path, str(temp_dir))

    image_path = get_images_dir(ppt_path)
    os.makedirs(image_path, exist_ok=True)

    pdf_path = temp_dir / f'{Path(ppt_path).stem}.pdf'
    convert_pdf_to_png(str(pdf_path), str(image_path))

    image_urls = encode_images_to_base64(image_path)

    return text_content, image_urls


def encode_image_to_base64(img_path):
    try:
        # Get file extension and corresponding MIME type
        extension = os.path.splitext(img_path)[1].lower()
        mime_types = {
            '.png': 'image/png',
            '.jpg': 'image/jpeg',
            '.jpeg': 'image/jpeg',
            '.gif': 'image/gif',
            '.bmp': 'image/bmp',
            '.webp': 'image/webp'
        }
        mime_type = mime_types.get(extension, 'image/jpeg')  # default to jpeg if unknown
        
        with open(img_path, 'rb') as f:
            image_data = f.read()
            base64_data = base64.b64encode(image_data).decode('utf-8')
            return f"data:{mime_type};base64,{base64_data}"
            
    except Exception as e:
        print(f"Error processing image {img_path}: {str(e)}")
        return None

def encode_images_to_base64(img_dir):
    image_files = os.listdir(img_dir)
    # Filter for files that match the pattern of "1.png", "2.png", etc.
    image_files = [f for f in image_files if re.match(r'^\d+\.png$', f)]
    image_files = sorted(image_files, key=lambda x: int(x.split("/")[-1].split('.')[0]))
    image_urls = []
    
    for filename in image_files:
        img_path = os.path.join(img_dir, filename)
        image_url = encode_image_to_base64(img_path)
        if image_url:
            image_urls.append(image_url)
    
    return image_urls


def extract_text_from_ppt(ppt_path: str) -> str:

    presentation = Presentation(ppt_path)

    all_content = {}
    for slide_number, slide in enumerate(presentation.slides):
        content = ""
        for shape in slide.shapes:

            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    content += paragraph.text + "\n"
                content += "\n"
        all_content[str(slide_number)] = content

    return all_content


def extract_json(text):
    """
    Extracts and parses JSON from a text string that may contain other content.
    
    This function attempts multiple strategies to extract valid JSON:
    1. First tries to find complete JSON structures using regex patterns
    2. Falls back to finding outermost { } or [ ] if regex fails
    3. Cleans the text to fix common JSON formatting issues
    4. Attempts multiple parsing methods with increasing aggressiveness
    
    Args:
        text: Input text that may contain JSON
        
    Returns:
        Parsed JSON as a dictionary or list
        
    Raises:
        ValueError: If JSON cannot be parsed after all attempts
    """
    import json
    import re
    
    # Clean up the text first to handle common issues
    # Remove markdown code block markers if present
    text = re.sub(r'```(?:json)?|```', '', text)
    
    # Strategy 1: Use regex to find JSON structures
    # Try to find a complete JSON object with balanced braces
    json_pattern = r'(\{(?:[^{}]|(?R))*\})'
    array_pattern = r'(\[(?:[^\[\]]|(?R))*\])'
    
    # Since Python's re doesn't support recursion (?R), we'll use a simpler approach
    # Try to match the outermost JSON object or array
    obj_match = re.search(r'\{.*\}', text, re.DOTALL)
    arr_match = re.search(r'\[.*\]', text, re.DOTALL)
    
    json_candidates = []
    if obj_match:
        json_candidates.append(obj_match.group(0))
    if arr_match:
        json_candidates.append(arr_match.group(0))
    
    # If we found potential JSON structures via regex
    for json_str in json_candidates:
        # Try to parse it directly first
        try:
            return json.loads(json_str)
        except json.JSONDecodeError:
            # If direct parsing fails, we'll continue to more aggressive methods
            pass
    
    # Strategy 2: Fall back to the original method if regex didn't work
    # Find leftmost { or [
    left_curly = text.find('{')
    left_bracket = text.find('[')
    
    # Determine which comes first (if both exist)
    if left_curly != -1 and left_bracket != -1:
        left_pos = min(left_curly, left_bracket)
    elif left_curly != -1:
        left_pos = left_curly
    elif left_bracket != -1:
        left_pos = left_bracket
    else:
        raise ValueError("No JSON structure found in text")
    
    # Find rightmost } or ]
    right_curly = text.rfind('}')
    right_bracket = text.rfind(']')
    
    # Determine which comes last (if both exist)
    if right_curly != -1 and right_bracket != -1:
        right_pos = max(right_curly, right_bracket)
    elif right_curly != -1:
        right_pos = right_curly
    elif right_bracket != -1:
        right_pos = right_bracket
    else:
        raise ValueError("No JSON structure found in text")
    
    # Extract the JSON substring
    json_str = text[left_pos:right_pos+1]
    
    # Strategy 3: Clean up common issues before parsing
    # Fix line breaks after commas which are common in LLM outputs
    json_str = re.sub(r',\s*\n\s*', ', ', json_str)
    
    # Fix missing quotes around keys
    json_str = re.sub(r'(\{|\,)\s*([a-zA-Z0-9_]+)\s*:', r'\1 "\2":', json_str)
    
    # Remove any trailing commas before closing brackets or braces
    json_str = re.sub(r',\s*(\}|\])', r'\1', json_str)
    
    # Strategy 4: Try multiple parsing methods with increasing aggressiveness
    # First attempt: standard json.loads
    try:
        return json.loads(json_str)
    except json.JSONDecodeError as e:
        print(f"Standard JSON parsing failed: {e}")
        
        # Second attempt: Try with ast.literal_eval for more forgiving parsing
        try:
            import ast
            return ast.literal_eval(json_str)
        except (SyntaxError, ValueError) as e:
            print(f"AST literal_eval failed: {e}")
            
            # Third attempt: Try with json5 if available (most lenient JSON parser)
            try:
                try:
                    import json5
                except ImportError:
                    # If json5 is not installed, try to install it
                    import subprocess
                    import sys
                    print("json5 module not found, attempting to install...")
                    subprocess.check_call([sys.executable, "-m", "pip", "install", "json5"])
                    import json5
                
                return json5.loads(json_str)
            except Exception as e:
                print(f"JSON5 parsing failed: {e}")
                
                # Final fallback: raise a more detailed error
                raise ValueError(f"Failed to parse JSON after multiple attempts. Original text: {text[:100]}...")


def convert_pdf_to_png(input_file, output_dir):
    os.makedirs(output_dir, exist_ok=True)

    # Open the PDF file
    doc = fitz.open(input_file)

    # Iterate through each page of the PDF
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)  # Load the current page
        pix = page.get_pixmap()  # Render page to an image
        output_path = os.path.join(output_dir, f'{page_num + 1}.png')  # Define output path for the image
        pix.save(output_path)  # Save the image

    print(f"Conversion completed. Images are saved in '{output_dir}'.")


def pptx_to_pdf(input_file: str, output_dir: str) -> bool:
    """Convert a PPTX file to PDF using LibreOffice in Docker.
    
    Args:
        input_file: Path to the PPTX file
        output_dir: Directory to save the PDF file
        
    Returns:
        bool: True if conversion successful, False otherwise
    """
    try:
        # Create output directory if it doesn't exist
        os.makedirs(output_dir, exist_ok=True)
        
        # Convert paths to absolute paths
        input_file = os.path.abspath(input_file)
        input_dir = os.path.dirname(input_file)
        
        # Run LibreOffice conversion in Docker
        cmd = f'docker run --rm -v "{input_dir}:/data" -v "{output_dir}:/output" "libreoffice-converter" libreoffice --headless --convert-to pdf --outdir /output "{os.path.basename(input_file)}"'
        
        # Execute command and check return code
        result = os.system(cmd)
        if result != 0:
            print(f"Error converting {input_file} to PDF")
            return False
            
        return True
        
    except Exception as e:
        print(f"Error during PPTX to PDF conversion: {str(e)}")
        return False


def extract_and_save_slides(accumulated_response, task_id, storage, saved_slides=None, is_final=False):
    """
    Extract slide content from accumulated response and save to storage.
    
    Args:
        accumulated_response: The text response containing slide content with markers
        task_id: The unique task identifier for storage
        storage: An instance of ScriptStorage to save slides
        saved_slides: Set of slide numbers that have already been saved (or None to initialize)
        is_final: Whether this is the final processing (True) or streaming (False)
        
    Returns:
        Tuple of (saved_slides, scripts) where:
            - saved_slides is a set of slide numbers that have been saved
            - scripts is a dictionary of all extracted scripts with slide numbers as keys
    """
    # Initialize saved_slides if not provided
    if saved_slides is None:
        saved_slides = set()
    
    # Pattern to extract slide content between markers
    slide_pattern = r'===SLIDE\s+(\d+)===\s*(.*?)(?=\s*===SLIDE\s+\d+===|\s*$)'
    
    # Extract all slides from the response
    matches = re.findall(slide_pattern, accumulated_response, re.DOTALL)
    if not matches:
        return saved_slides, {}
    
    # Dictionary to store all scripts (for return value)
    scripts = {}
    
    # Determine which slides to process:
    # - During streaming (is_final=False): Process all except the last one (which might be incomplete)
    # - During final processing (is_final=True): Process all slides
    slides_to_process = matches if is_final else matches[:-1]
    
    # Add all slides to the scripts dictionary
    for slide_num, content in matches:
        slide_num_str = slide_num.strip()
        content = content.strip()
        scripts[slide_num_str] = content
    
    # Save slides that should be processed and haven't been saved yet
    for slide_num, content in slides_to_process:
        slide_num_int = int(slide_num.strip())
        content = content.strip()
        
        if slide_num_int not in saved_slides:
            storage.save_slide(task_id, slide_num_int, content)
            saved_slides.add(slide_num_int)
            logger.debug(f"Saved slide {slide_num_int} to storage")
    
    return saved_slides, scripts

