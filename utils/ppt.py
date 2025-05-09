"""
PowerPoint and PDF preprocessing utilities for text and image extraction.
"""

import os
import logging
from pptx import Presentation
from pathlib import Path
from typing import Dict, List, Optional, Tuple, Union

# Assuming PyMuPDF is installed for PDF text extraction
import fitz  # PyMuPDF

from utils.paths import get_temp_dir, get_images_dir
from utils.image import encode_images_to_base64, convert_pdf_to_png
# Assume utils.pdf might have more complex extraction, but for now use PyMuPDF directly
# from utils.pdf import extract_text_from_pdf

# Configure logging
logger = logging.getLogger(__name__)

def extract_text_from_pdf_simple(pdf_path: str) -> Dict[str, str]:
    """
    Extract text content from a PDF file, page by page.
    
    Args:
        pdf_path: Path to the PDF file
        
    Returns:
        dict: Dictionary mapping 0-indexed page numbers (as strings) to text content
    """
    all_content = {}
    try:
        doc = fitz.open(pdf_path)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            text = page.get_text("text")
            all_content[str(page_num)] = text.strip()
        doc.close()
    except Exception as e:
        logger.error(f"Error extracting text from PDF {pdf_path}: {e}", exc_info=True)
        # Return empty dict on error to avoid breaking callers expecting a dict
        return {}
    return all_content

def preprocess_ppt(input_path: Union[str, Path]) -> Tuple[Optional[Dict[str, str]], List[str]]:
    """
    Preprocess a PowerPoint (.pptx) or PDF (.pdf) file to extract text and convert slides/pages to images.
    
    Args:
        input_path: Path to the input file (.pptx or .pdf)
        
    Returns:
        tuple: Tuple containing (text_content, image_urls)
               - text_content: Dictionary mapping 0-indexed slide/page numbers to text, or None if extraction fails for PPTX.
               - image_urls: List of base64 encoded image data URLs.
    """
    if isinstance(input_path, str):
        input_path = Path(input_path)

    if not input_path.exists():
        raise FileNotFoundError(f"Input file not found: {input_path}")

    file_suffix = input_path.suffix.lower()
    image_urls = []
    text_content = None

    try:
        image_path_obj = get_images_dir(input_path)
        os.makedirs(image_path_obj, exist_ok=True)
        image_path_str = str(image_path_obj)

        if file_suffix == '.pptx':
            logger.info(f"Preprocessing PPTX: {input_path.name}")
            text_content = extract_text_from_ppt(str(input_path))
            
            # Conversion Path: PPTX -> PDF -> Images
            temp_dir_obj = get_temp_dir(input_path)
            os.makedirs(temp_dir_obj, exist_ok=True)
            temp_dir_str = str(temp_dir_obj)
            pdf_path_obj = temp_dir_obj / f'{input_path.stem}.pdf'
            pdf_path_str = str(pdf_path_obj)

            pdf_converted = pptx_to_pdf(str(input_path), temp_dir_str)
            if pdf_converted:
                convert_pdf_to_png(pdf_path_str, image_path_str)
                image_urls = encode_images_to_base64(image_path_obj)
            else:
                logger.error(f"Failed to convert PPTX to PDF for image generation: {input_path.name}")

        elif file_suffix == '.pdf':
            logger.info(f"Preprocessing PDF: {input_path.name}")
            # Extract text directly from PDF
            text_content = extract_text_from_pdf_simple(str(input_path))
            
            # Convert PDF pages directly to images
            convert_pdf_to_png(str(input_path), image_path_str)
            image_urls = encode_images_to_base64(image_path_obj)

        else:
            raise ValueError(f"Unsupported file type: {input_path}. Only .pptx and .pdf are supported.")

    except Exception as e:
        logger.error(f"Error during preprocessing {input_path.name}: {e}", exc_info=True)
        # Return None, empty list to indicate failure but allow callers to potentially handle it
        return None, []

    if not image_urls:
        logger.warning(f"Failed to generate images for {input_path.name}")
        # Text content might still be valid, so return it

    logger.info(f"Preprocessing complete for {input_path.name}. Text slides: {len(text_content) if text_content else 'N/A'}. Images: {len(image_urls)}.")
    return text_content, image_urls

def extract_text_from_ppt(ppt_path: str) -> dict:
    """
    Extract text content from a PowerPoint file.
    
    Args:
        ppt_path: Path to the PowerPoint file
        
    Returns:
        dict: Dictionary mapping 0-indexed slide numbers (as strings) to text content
    """
    all_content = {}
    try:
        presentation = Presentation(ppt_path)
        for slide_number, slide in enumerate(presentation.slides):
            content = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        content += paragraph.text + "\n"
                    # Add a separator between shapes, but only if shape had text
                    if shape.has_text_frame and shape.text_frame.text.strip():
                         content += "\n"
            # Store even if content is empty, key is 0-indexed slide number
            all_content[str(slide_number)] = content.strip()
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {ppt_path}: {e}", exc_info=True)
        return {}
    return all_content

def pptx_to_pdf(input_file: str, output_dir: str) -> bool:
    """
    Convert a PPTX file to PDF using LibreOffice in Docker.
    
    Args:
        input_file: Path to the PPTX file
        output_dir: Directory to save the PDF file
        
    Returns:
        bool: True if conversion successful, False otherwise
    """
    try:
        # Create output directory if it doesn't exist
        os.makedirs(output_dir, exist_ok=True)
        
        # Convert paths to absolute paths
        input_file = os.path.abspath(input_file)
        input_dir = os.path.dirname(input_file)
        
        # Run LibreOffice conversion in Docker
        cmd = f'docker run --rm -v "{input_dir}:/data" -v "{output_dir}:/output" "libreoffice-converter" libreoffice --headless --convert-to pdf --outdir /output "{os.path.basename(input_file)}"'
        
        # Execute command and check return code
        result = os.system(cmd)
        if result != 0:
            logger.error(f"Error converting {input_file} to PDF")
            return False
            
        return True
        
    except Exception as e:
        logger.error(f"Error during PPTX to PDF conversion: {str(e)}")
        return False 